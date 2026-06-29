import os
import re
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize Presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Scheme
COLOR_PRIMARY = RGBColor(0, 102, 204)       # Action Blue (#0066cc)
COLOR_TEXT_MAIN = RGBColor(29, 29, 31)     # Ink (#1d1d1f)
COLOR_TEXT_MUTED = RGBColor(120, 120, 120)  # Gray
COLOR_BG_BODY = RGBColor(245, 245, 247)     # Light Gray
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(197, 34, 31)        # Accent Red for warnings
COLOR_CALLOUT_BG = RGBColor(230, 242, 255)  # Light Blue Callout

# Load HTML file
html_path = r"c:\bigdata_vibe1\uti_slides.html"
with open(html_path, "r", encoding="utf-8") as f:
    soup = BeautifulSoup(f.read(), "html.parser")

slides = soup.find_all("section", class_="slide-card")

def set_background_color(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title_text, subtitle_text, meta_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    set_background_color(slide, COLOR_PRIMARY)
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text.replace("<br>", "\n")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT
    
    # Subtitle Box
    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.333), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = subtitle_text
    p_sub.font.size = Pt(22)
    p_sub.font.color.rgb = COLOR_WHITE
    p_sub.alignment = PP_ALIGN.LEFT
    
    # Meta Box
    meta_box = slide.shapes.add_textbox(Inches(1.0), Inches(5.8), Inches(11.333), Inches(1.0))
    tf_meta = meta_box.text_frame
    tf_meta.word_wrap = True
    p_meta = tf_meta.paragraphs[0]
    p_meta.text = meta_text.strip()
    p_meta.font.size = Pt(14)
    p_meta.font.color.rgb = COLOR_WHITE
    p_meta.alignment = PP_ALIGN.LEFT

def format_bullet_point(tf, text, is_first=False, is_bold=False):
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    p.space_after = Pt(8)
    p.level = 0
    
    # Bold check
    if is_bold:
        p.font.bold = True
        p.font.size = Pt(16)
        p.text = text
    else:
        # Check for inline strong tags or text formatting
        p.text = text
        p.font.size = Pt(16)
    
    p.font.color.rgb = COLOR_TEXT_MAIN

def parse_html_text(element):
    if not element:
        return ""
    # Replace strong/em tags with plain text for simplicity or keep tags
    return element.get_text(" ", strip=True)

# Loop through all slides
for i, slide_elem in enumerate(slides):
    # Check if it's title slide
    if "title-slide" in slide_elem.get("class", []):
        title = slide_elem.find(class_="slide-title")
        subtitle = slide_elem.find(class_="slide-subtitle")
        meta = slide_elem.find(class_="meta-info")
        
        title_txt = parse_html_text(title) if title else "UTI 임상진료지침 권고안"
        sub_txt = parse_html_text(subtitle) if subtitle else ""
        meta_txt = parse_html_text(meta) if meta else ""
        
        add_title_slide(prs, title_txt, sub_txt, meta_txt)
        continue

    # Regular slide creation
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    set_background_color(slide, COLOR_WHITE)
    
    # Extract Category & Title
    category_elem = slide_elem.find(class_="slide-category")
    title_elem = slide_elem.find(class_="slide-title")
    category_txt = parse_html_text(category_elem).upper() if category_elem else ""
    title_txt = parse_html_text(title_elem) if title_elem else ""
    
    # 1. Draw Title Area
    # Category Text
    if category_txt:
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
        p = cat_box.text_frame.paragraphs[0]
        p.text = category_txt
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRIMARY
        
    # Title Text
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    p = title_tf.paragraphs[0]
    p.text = title_txt
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    # Add a thin separator line under title
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.color.rgb = COLOR_PRIMARY
    
    # Extract Body
    body_elem = slide_elem.find(class_="slide-body")
    if not body_elem:
        continue
        
    # Positioning for main contents
    content_top = Inches(1.8)
    content_height = Inches(5.0)
    
    # Check for grid columns
    grid_2 = body_elem.find(class_="grid-2")
    callout = body_elem.find(class_="callout-box")
    table = body_elem.find("table")
    
    if table:
        # Create Table inside PPTX
        rows_data = []
        headers = [parse_html_text(th) for th in table.find_all("th")]
        if headers:
            rows_data.append(headers)
        
        for tr in table.find_all("tr"):
            cells = [parse_html_text(td) for td in tr.find_all("td")]
            if cells:
                rows_data.append(cells)
                
        if rows_data:
            num_rows = len(rows_data)
            num_cols = len(rows_data[0])
            
            # Position Table
            t_left = Inches(0.8)
            t_top = Inches(2.2)
            t_width = Inches(11.7)
            t_height = Inches(0.4 * num_rows)
            
            shape = slide.shapes.add_table(num_rows, num_cols, t_left, t_top, t_width, t_height)
            pptx_table = shape.table
            
            col_width = int(t_width / num_cols)
            for col_idx in range(num_cols):
                pptx_table.columns[col_idx].width = col_width
                
            # Fill cells
            for row_idx, row_data in enumerate(rows_data):
                for col_idx, cell_value in enumerate(row_data):
                    cell = pptx_table.cell(row_idx, col_idx)
                    cell.text = cell_value
                    # Format cell text
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(13)
                        paragraph.font.color.rgb = COLOR_TEXT_MAIN
                        if row_idx == 0:
                            paragraph.font.bold = True
                            paragraph.font.color.rgb = COLOR_WHITE
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = COLOR_PRIMARY
                            
    elif grid_2:
        # Two columns layout
        col_elements = grid_2.find_all(recursive=False)
        if len(col_elements) >= 2:
            left_col = col_elements[0]
            right_col = col_elements[1]
            
            # Left box
            left_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(5.6), content_height)
            left_tf = left_box.text_frame
            left_tf.word_wrap = True
            
            # Right box
            right_box = slide.shapes.add_textbox(Inches(6.8), content_top, Inches(5.7), content_height)
            right_tf = right_box.text_frame
            right_tf.word_wrap = True
            
            # Populate left col
            left_items = left_col.find_all(["p", "li", "h3", "h4"])
            for idx, item in enumerate(left_items):
                text = parse_html_text(item)
                format_bullet_point(left_tf, text, is_first=(idx==0), is_bold=(item.name in ["h3", "h4"]))
                
            # Populate right col
            right_items = right_col.find_all(["p", "li", "h3", "h4"])
            for idx, item in enumerate(right_items):
                text = parse_html_text(item)
                format_bullet_point(right_tf, text, is_first=(idx==0), is_bold=(item.name in ["h3", "h4"]))
    else:
        # Single column layout (with or without Callout)
        text_box = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11.7), content_height)
        tf = text_box.text_frame
        tf.word_wrap = True
        
        # If callout exists, render it
        if callout:
            callout_title = callout.find(class_="callout-box-title")
            title_text = parse_html_text(callout_title) if callout_title else "참고"
            
            # Remove title node to get clean body text
            if callout_title:
                callout_title.decompose()
            body_text = parse_html_text(callout)
            
            # Draw Callout Background Box
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.2))
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLOR_CALLOUT_BG
            shape.line.color.rgb = COLOR_PRIMARY
            
            callout_tf = shape.text_frame
            callout_tf.word_wrap = True
            p1 = callout_tf.paragraphs[0]
            p1.text = f"💡 {title_text}: {body_text}"
            p1.font.size = Pt(14)
            p1.font.color.rgb = COLOR_PRIMARY
            p1.font.bold = True
            
            # Adjust main text box top position so it doesn't overlap the callout
            text_box.top = Inches(3.3)
            text_box.height = Inches(3.5)
            
        # Add normal items
        items = body_elem.find_all(["p", "li", "h3", "h4"])
        # Filter items that were inside callout (they might have been decomposed, but safety check)
        items = [i for i in items if not i.find_parent(class_="callout-box") and i.name != "th" and i.name != "td"]
        
        for idx, item in enumerate(items):
            text = parse_html_text(item)
            if not text:
                continue
            format_bullet_point(tf, text, is_first=(idx==0), is_bold=(item.name in ["h3", "h4"]))

# Save Presentation
output_pptx = r"c:\bigdata_vibe1\uti_slides.pptx"
prs.save(output_pptx)
print(f"Successfully converted HTML slides to: {output_pptx}")
